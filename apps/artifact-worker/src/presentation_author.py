from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.util import Inches, Pt

from .text_runs import replace_text_in_runs

ACCENT = RGBColor(43, 91, 175)
TEXT = RGBColor(23, 32, 51)
MUTED = RGBColor(91, 104, 128)


def _set_run_style(
    run, *, size: int, bold: bool = False, color: RGBColor = TEXT
) -> None:
    run.font.name = "Liberation Sans"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int,
    bold: bool = False,
    color: RGBColor = TEXT,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    _set_run_style(paragraph.runs[0], size=size, bold=bold, color=color)
    return box


def _bullet_font_size(items: list[str], two_column: bool = False) -> int:
    longest = max((len(item) for item in items), default=0)
    total = sum(len(item) for item in items)
    budget = 450 if two_column else 750
    if longest > 220 or total > budget:
        return 18
    if longest > 140 or total > budget * 0.72:
        return 21
    return 24


def _add_bullets(
    slide,
    items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    two_column: bool = False,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    size = _bullet_font_size(items, two_column)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        _set_run_style(paragraph.runs[0], size=size)


def _add_slide(presentation: PresentationType, spec: dict[str, Any]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    layout = spec["layout"]
    if layout == "title":
        title = _add_textbox(
            slide, 0.8, 2.2, 11.8, 1.4, spec["title"], size=40, bold=True
        )
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if spec.get("subtitle"):
            subtitle = _add_textbox(
                slide, 1.4, 3.8, 10.6, 0.8, spec["subtitle"], size=22, color=MUTED
            )
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    elif layout == "section":
        _add_textbox(
            slide, 0.8, 2.5, 11.6, 1.3, spec["title"], size=36, bold=True, color=ACCENT
        )
        if spec.get("subtitle"):
            _add_textbox(
                slide, 0.8, 3.9, 11.2, 0.8, spec["subtitle"], size=22, color=MUTED
            )
    elif layout == "two_column":
        _add_textbox(slide, 0.7, 0.45, 12, 0.8, spec["title"], size=30, bold=True)
        _add_bullets(
            slide, spec.get("leftBullets", []), 0.7, 1.5, 5.8, 4.9, two_column=True
        )
        _add_bullets(
            slide, spec.get("rightBullets", []), 6.8, 1.5, 5.8, 4.9, two_column=True
        )
    else:
        _add_textbox(slide, 0.7, 0.45, 12, 0.8, spec["title"], size=30, bold=True)
        _add_bullets(slide, spec.get("bullets", []), 0.9, 1.5, 11.5, 4.9)
    if spec.get("footer"):
        _add_textbox(slide, 0.75, 7.05, 11.8, 0.25, spec["footer"], size=9, color=MUTED)
    if spec.get("speakerNotes"):
        slide.notes_slide.notes_text_frame.text = spec["speakerNotes"]


def create_presentation(spec: dict[str, Any], destination: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = spec["title"]
    if spec.get("author"):
        presentation.core_properties.author = spec["author"]
    for slide_spec in spec["slides"]:
        _add_slide(presentation, slide_spec)
    presentation.save(destination)


def _replace_text(
    presentation: PresentationType, find: str, replace: str, replace_all: bool
) -> int:
    count = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                replacements = replace_text_in_runs(
                    paragraph.runs,
                    find,
                    replace,
                    replace_all,
                )
                count += replacements
                if replacements > 0 and not replace_all:
                    return count
    return count


def _delete_slide(presentation: PresentationType, slide_number: int) -> None:
    slide_id_list = presentation.slides._sldIdLst  # noqa: SLF001
    slide_id = slide_id_list[slide_number - 1]
    presentation.part.drop_rel(slide_id.rId)
    slide_id_list.remove(slide_id)


def _move_last_slide_after(
    presentation: PresentationType,
    after_slide: int,
) -> None:
    slide_id_list = presentation.slides._sldIdLst  # noqa: SLF001
    slide_id = slide_id_list[-1]
    slide_id_list.remove(slide_id)
    slide_id_list.insert(after_slide, slide_id)


def edit_presentation(
    source: Path, edits: list[dict[str, Any]], destination: Path
) -> None:
    presentation = Presentation(source)
    for edit in edits:
        edit_type = edit["type"]
        if edit_type == "insert_slides":
            after_slide = edit["afterSlide"]
            if after_slide > len(presentation.slides):
                raise ValueError("Presentation insertion position does not exist.")
            for slide_spec in edit["slides"]:
                _add_slide(presentation, slide_spec)
                _move_last_slide_after(presentation, after_slide)
                after_slide += 1
        elif edit_type == "replace_text":
            replacements = _replace_text(
                presentation,
                edit["find"],
                edit["replace"],
                edit.get("replaceAll", False),
            )
            if replacements == 0:
                raise ValueError(f"Presentation text was not found: {edit['find']}")
        elif edit_type == "delete_slides":
            for slide_number in sorted(set(edit["slideNumbers"]), reverse=True):
                if slide_number > len(presentation.slides):
                    raise ValueError(
                        f"Presentation slide was not found: {slide_number}"
                    )
                _delete_slide(presentation, slide_number)
        else:
            raise ValueError(f"Unsupported presentation edit: {edit_type}")
    if len(presentation.slides) == 0:
        raise ValueError("A presentation must contain at least one slide.")
    presentation.save(destination)
