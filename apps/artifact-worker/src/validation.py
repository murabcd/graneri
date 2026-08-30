from __future__ import annotations

from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import BadZipFile, ZipFile

from docx import Document
from openpyxl import load_workbook
from PIL import Image, ImageStat
from pptx import Presentation
from pypdf import PdfReader

from .process import convert_with_libreoffice, render_pdf

FORMULA_ERROR_VALUES = {
    "#DIV/0!",
    "#N/A",
    "#NAME?",
    "#NULL!",
    "#NUM!",
    "#REF!",
    "#VALUE!",
}
MAX_PDF_PAGES = 200
MAX_PRESENTATION_SLIDES = 100
MAX_SPREADSHEET_CELLS = 250_000
MAX_SPREADSHEET_SHEETS = 20


def _assert_ooxml(path: Path, required_prefix: str) -> None:
    try:
        with ZipFile(path) as archive:
            names = archive.namelist()
            if "[Content_Types].xml" not in names:
                raise ValueError("OOXML content types are missing.")
            if not any(name.startswith(required_prefix) for name in names):
                raise ValueError(f"OOXML package is missing {required_prefix} content.")
            corrupt = archive.testzip()
            if corrupt:
                raise ValueError(f"OOXML package contains a corrupt member: {corrupt}")
    except BadZipFile as error:
        raise ValueError("Artifact is not a valid OOXML package.") from error


def _validate_docx(path: Path) -> None:
    _assert_ooxml(path, "word/")
    document = Document(path)
    if not document.paragraphs and not document.tables:
        raise ValueError("DOCX contains no document content.")


def _assert_no_formula_errors(path: Path, *, data_only: bool) -> None:
    workbook = load_workbook(path, data_only=data_only, read_only=True)
    try:
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows():
                for cell in row:
                    if (
                        isinstance(cell.value, str)
                        and cell.value in FORMULA_ERROR_VALUES
                    ):
                        raise ValueError(
                            f"XLSX contains formula error {cell.value} at "
                            f"{worksheet.title}!{cell.coordinate}."
                        )
    finally:
        workbook.close()


def _assert_xlsx_bounds(path: Path) -> None:
    workbook = load_workbook(path, data_only=False, read_only=True)
    try:
        if not workbook.sheetnames:
            raise ValueError("XLSX contains no sheets.")
        if len(workbook.sheetnames) > MAX_SPREADSHEET_SHEETS:
            raise ValueError(
                f"XLSX exceeds the {MAX_SPREADSHEET_SHEETS}-sheet validation limit."
            )
        cell_count = 0
        for worksheet in workbook.worksheets:
            if worksheet.max_row < 1 or worksheet.max_column < 1:
                raise ValueError(f"XLSX sheet is empty: {worksheet.title}")
            cell_count += worksheet.max_row * worksheet.max_column
            if cell_count > MAX_SPREADSHEET_CELLS:
                raise ValueError(
                    f"XLSX exceeds the {MAX_SPREADSHEET_CELLS:,}-cell validation limit."
                )
    finally:
        workbook.close()


def _validate_xlsx(path: Path, validation_directory: Path) -> None:
    _assert_ooxml(path, "xl/")
    _assert_xlsx_bounds(path)
    _assert_no_formula_errors(path, data_only=False)
    recalculated_path = convert_with_libreoffice(
        path,
        validation_directory / "recalculated",
        "xlsx",
    )
    _assert_no_formula_errors(recalculated_path, data_only=True)


def _validate_pptx(path: Path) -> None:
    _assert_ooxml(path, "ppt/")
    presentation = Presentation(path)
    if len(presentation.slides) == 0:
        raise ValueError("PPTX contains no slides.")
    if len(presentation.slides) > MAX_PRESENTATION_SLIDES:
        raise ValueError(
            f"PPTX exceeds the {MAX_PRESENTATION_SLIDES}-slide validation limit."
        )


def _validate_pdf(path: Path) -> int:
    reader = PdfReader(path)
    if len(reader.pages) == 0:
        raise ValueError("PDF contains no pages.")
    if len(reader.pages) > MAX_PDF_PAGES:
        raise ValueError(f"PDF exceeds the {MAX_PDF_PAGES}-page validation limit.")
    for index, page in enumerate(reader.pages, start=1):
        box = page.mediabox
        if float(box.width) <= 0 or float(box.height) <= 0:
            raise ValueError(f"PDF page {index} has invalid dimensions.")
    return len(reader.pages)


def _assert_render_not_blank(image_path: Path) -> None:
    with Image.open(image_path) as image:
        image.verify()
    with Image.open(image_path).convert("RGB") as image:
        if image.width < 100 or image.height < 100:
            raise ValueError(f"Rendered page is unexpectedly small: {image_path.name}")
        thumbnail = image.copy()
        thumbnail.thumbnail((256, 256))
        statistics = ImageStat.Stat(thumbnail)
        if all(mean > 254.8 for mean in statistics.mean) and all(
            deviation < 0.6 for deviation in statistics.stddev
        ):
            raise ValueError(f"Rendered page appears blank: {image_path.name}")


def _render_and_validate(pdf_path: Path, expected_minimum_pages: int) -> None:
    with TemporaryDirectory(prefix="artifact-render-") as directory:
        render_prefix = Path(directory) / "page"
        render_pdf(pdf_path, render_prefix)
        images = sorted(Path(directory).glob("page-*.png"))
        if len(images) < expected_minimum_pages:
            raise ValueError("Rendered artifact is missing pages or slides.")
        for image in images:
            _assert_render_not_blank(image)


def validate_source_artifact(path: Path, format_name: str) -> None:
    if not path.is_file() or path.stat().st_size == 0:
        raise ValueError("Artifact source is empty.")
    if format_name == "docx":
        _validate_docx(path)
    elif format_name == "xlsx":
        _assert_ooxml(path, "xl/")
        _assert_xlsx_bounds(path)
    elif format_name == "pptx":
        _validate_pptx(path)
    elif format_name == "pdf":
        _validate_pdf(path)
    else:
        raise ValueError(f"Unsupported artifact source format: {format_name}")


def validate_artifact(path: Path, format_name: str) -> None:
    if not path.is_file() or path.stat().st_size == 0:
        raise ValueError("Artifact output is empty.")
    with TemporaryDirectory(prefix="artifact-validation-") as directory:
        validation_directory = Path(directory)
        if format_name == "docx":
            _validate_docx(path)
            pdf_path = convert_with_libreoffice(path, validation_directory, "pdf")
            page_count = _validate_pdf(pdf_path)
            _render_and_validate(pdf_path, page_count)
        elif format_name == "xlsx":
            _validate_xlsx(path, validation_directory)
            pdf_path = convert_with_libreoffice(path, validation_directory, "pdf")
            page_count = _validate_pdf(pdf_path)
            _render_and_validate(pdf_path, page_count)
        elif format_name == "pptx":
            _validate_pptx(path)
            slide_count = len(Presentation(path).slides)
            pdf_path = convert_with_libreoffice(path, validation_directory, "pdf")
            page_count = _validate_pdf(pdf_path)
            if page_count != slide_count:
                raise ValueError(
                    "Rendered presentation page count does not match its slide count."
                )
            _render_and_validate(pdf_path, slide_count)
        elif format_name == "pdf":
            page_count = _validate_pdf(path)
            _render_and_validate(path, page_count)
        else:
            raise ValueError(f"Unsupported artifact validation format: {format_name}")
